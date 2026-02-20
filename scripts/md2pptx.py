#!/usr/bin/env python3
"""Convert a markdown slide deck (--- separated) into an editable .pptx file.

Usage:
    python md2pptx.py input.md -o output.pptx [--title "Deck Title"]

Markdown format:
    # Slide Title        -> title slide (first slide) or section header
    ## Slide Heading     -> content slide heading
    - bullet             -> bulleted list
    1. item              -> numbered list
    ```code```           -> code block (monospace, gray background)
    > quote              -> italic blockquote
    **bold**, *italic*   -> inline formatting
    | col | col |        -> table
    plain text           -> body paragraph
    ---                  -> slide separator
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Theme colours ──────────────────────────────────────────────────────────
COLOR_HEADING = RGBColor(0x1E, 0x3A, 0x5F)
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_CODE_BG = RGBColor(0xF0, 0xF0, 0xF0)
COLOR_CODE_FG = RGBColor(0x2D, 0x2D, 0x2D)
COLOR_QUOTE = RGBColor(0x55, 0x55, 0x55)
COLOR_ACCENT = RGBColor(0x25, 0x63, 0xEB)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TABLE_HEADER_BG = RGBColor(0x25, 0x63, 0xEB)
COLOR_TABLE_ALT_BG = RGBColor(0xF5, 0xF7, 0xFA)

FONT_SANS = "Calibri"
FONT_MONO = "Consolas"

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Content area
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP_TITLE = Inches(0.6)
MARGIN_TOP_BODY = Inches(1.8)
CONTENT_WIDTH = Inches(11.7)
CONTENT_HEIGHT_FULL = Inches(5.8)
CONTENT_HEIGHT_WITH_TITLE = Inches(5.0)
TITLE_HEIGHT = Inches(1.0)


# ── Inline formatting ─────────────────────────────────────────────────────
def _add_formatted_runs(paragraph, text: str, base_size: int = 18,
                        base_color: RGBColor = COLOR_BODY,
                        base_bold: bool = False, base_italic: bool = False,
                        base_font: str = FONT_SANS):
    """Parse markdown inline formatting (**bold**, *italic*, `code`) and add runs."""
    # Pattern matches **bold**, *italic*, `code`, or plain text
    pattern = re.compile(r'(\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`|([^*`]+))')
    for match in pattern.finditer(text):
        if match.group(2) is not None:  # **bold**
            run = paragraph.add_run()
            run.text = match.group(2)
            run.font.size = Pt(base_size)
            run.font.color.rgb = base_color
            run.font.bold = True
            run.font.italic = base_italic
            run.font.name = base_font
        elif match.group(3) is not None:  # *italic*
            run = paragraph.add_run()
            run.text = match.group(3)
            run.font.size = Pt(base_size)
            run.font.color.rgb = base_color
            run.font.bold = base_bold
            run.font.italic = True
            run.font.name = base_font
        elif match.group(4) is not None:  # `code`
            run = paragraph.add_run()
            run.text = match.group(4)
            run.font.size = Pt(base_size - 2)
            run.font.color.rgb = COLOR_CODE_FG
            run.font.bold = False
            run.font.italic = False
            run.font.name = FONT_MONO
        elif match.group(5) is not None:  # plain text
            run = paragraph.add_run()
            run.text = match.group(5)
            run.font.size = Pt(base_size)
            run.font.color.rgb = base_color
            run.font.bold = base_bold
            run.font.italic = base_italic
            run.font.name = base_font


# ── Slide content parsing ─────────────────────────────────────────────────
def _parse_slide_blocks(md: str) -> list[dict]:
    """Parse a single slide's markdown into structured blocks."""
    blocks: list[dict] = []
    lines = md.split("\n")
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Skip empty lines
        if not stripped:
            i += 1
            continue

        # Heading
        if stripped.startswith("# ") and not stripped.startswith("## "):
            blocks.append({"type": "h1", "text": stripped[2:].strip()})
            i += 1
            continue

        if stripped.startswith("## "):
            blocks.append({"type": "h2", "text": stripped[3:].strip()})
            i += 1
            continue

        if stripped.startswith("### "):
            blocks.append({"type": "h3", "text": stripped[4:].strip()})
            i += 1
            continue

        # Code block
        if stripped.startswith("```"):
            code_lines: list[str] = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing ```
            blocks.append({"type": "code", "text": "\n".join(code_lines)})
            continue

        # Table
        if "|" in stripped and stripped.startswith("|"):
            table_lines: list[str] = []
            while i < len(lines) and "|" in lines[i].strip():
                table_lines.append(lines[i].strip())
                i += 1
            blocks.append({"type": "table", "lines": table_lines})
            continue

        # Blockquote
        if stripped.startswith("> "):
            quote_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("> "):
                quote_lines.append(lines[i].strip()[2:])
                i += 1
            blocks.append({"type": "quote", "text": " ".join(quote_lines)})
            continue

        # Bullet list
        if re.match(r'^[-*]\s', stripped):
            items: list[str] = []
            while i < len(lines) and re.match(r'^\s*[-*]\s', lines[i]):
                items.append(re.sub(r'^\s*[-*]\s+', '', lines[i]).strip())
                i += 1
            blocks.append({"type": "bullets", "items": items})
            continue

        # Numbered list
        if re.match(r'^\d+\.\s', stripped):
            items = []
            while i < len(lines) and re.match(r'^\s*\d+\.\s', lines[i]):
                items.append(re.sub(r'^\s*\d+\.\s+', '', lines[i]).strip())
                i += 1
            blocks.append({"type": "numbered", "items": items})
            continue

        # Plain paragraph
        para_lines: list[str] = []
        while i < len(lines) and lines[i].strip() and not any([
            lines[i].strip().startswith("#"),
            lines[i].strip().startswith("```"),
            lines[i].strip().startswith("> "),
            lines[i].strip().startswith("| "),
            re.match(r'^\s*[-*]\s', lines[i]),
            re.match(r'^\s*\d+\.\s', lines[i]),
        ]):
            para_lines.append(lines[i].strip())
            i += 1
        if para_lines:
            blocks.append({"type": "paragraph", "text": " ".join(para_lines)})

    return blocks


def _is_title_slide(blocks: list[dict]) -> bool:
    """Check if this looks like a title-only slide (just h1 + optional paragraph)."""
    non_empty = [b for b in blocks if b["type"] != "paragraph" or b.get("text", "").strip()]
    types = [b["type"] for b in non_empty]
    if types == ["h1"]:
        return True
    if types == ["h1", "paragraph"]:
        return True
    return False


# ── Slide rendering ───────────────────────────────────────────────────────
def _render_title_slide(prs: Presentation, blocks: list[dict]):
    """Render a title/cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background - accent color
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_ACCENT

    # Title text
    h1_text = blocks[0]["text"]
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_formatted_runs(p, h1_text, base_size=40, base_color=COLOR_WHITE,
                        base_bold=True)

    # Subtitle
    subtitle_blocks = [b for b in blocks[1:] if b["type"] == "paragraph"]
    if subtitle_blocks:
        txBox2 = slide.shapes.add_textbox(
            Inches(2.0), Inches(4.4), Inches(9.3), Inches(1.2)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        _add_formatted_runs(p2, subtitle_blocks[0]["text"], base_size=20,
                            base_color=COLOR_WHITE, base_italic=True)


def _render_content_slide(prs: Presentation, blocks: list[dict]):
    """Render a standard content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Determine if there's a heading
    heading_text = None
    content_blocks = blocks
    body_top = MARGIN_TOP_BODY
    body_height = CONTENT_HEIGHT_WITH_TITLE

    if blocks and blocks[0]["type"] in ("h1", "h2", "h3"):
        heading_text = blocks[0]["text"]
        content_blocks = blocks[1:]
    else:
        body_top = Inches(0.8)
        body_height = CONTENT_HEIGHT_FULL

    # Add heading
    if heading_text:
        txBox = slide.shapes.add_textbox(
            MARGIN_LEFT, MARGIN_TOP_TITLE, CONTENT_WIDTH, TITLE_HEIGHT
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_formatted_runs(p, heading_text, base_size=32, base_color=COLOR_HEADING,
                            base_bold=True)

        # Add accent underline
        line_shape = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            MARGIN_LEFT, Inches(1.55), Inches(2.0), Pt(3)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = COLOR_ACCENT
        line_shape.line.fill.background()

    # Render content blocks
    if not content_blocks:
        return

    # For simple content (just bullets or paragraphs), use a single text box
    current_top = body_top

    for block in content_blocks:
        if block["type"] == "bullets":
            txBox = slide.shapes.add_textbox(
                MARGIN_LEFT, current_top, CONTENT_WIDTH,
                Inches(min(len(block["items"]) * 0.45, float(body_height / Inches(1))))
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, item in enumerate(block["items"]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_before = Pt(6)
                p.space_after = Pt(4)
                # Bullet character
                run_bullet = p.add_run()
                run_bullet.text = "\u2022  "
                run_bullet.font.size = Pt(18)
                run_bullet.font.color.rgb = COLOR_ACCENT
                run_bullet.font.name = FONT_SANS
                _add_formatted_runs(p, item, base_size=18)
            current_top += Inches(len(block["items"]) * 0.45)

        elif block["type"] == "numbered":
            txBox = slide.shapes.add_textbox(
                MARGIN_LEFT, current_top, CONTENT_WIDTH,
                Inches(min(len(block["items"]) * 0.45, float(body_height / Inches(1))))
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, item in enumerate(block["items"]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_before = Pt(6)
                p.space_after = Pt(4)
                run_num = p.add_run()
                run_num.text = f"{j + 1}.  "
                run_num.font.size = Pt(18)
                run_num.font.color.rgb = COLOR_ACCENT
                run_num.font.bold = True
                run_num.font.name = FONT_SANS
                _add_formatted_runs(p, item, base_size=18)
            current_top += Inches(len(block["items"]) * 0.45)

        elif block["type"] == "code":
            code_height = Inches(min(len(block["text"].split("\n")) * 0.3 + 0.4, 4.0))
            shape = slide.shapes.add_shape(
                1,  # rectangle
                MARGIN_LEFT, current_top, CONTENT_WIDTH, code_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CODE_BG
            shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            shape.line.width = Pt(1)

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(14)
            tf.margin_right = Pt(14)
            tf.margin_top = Pt(10)
            tf.margin_bottom = Pt(10)

            for j, code_line in enumerate(block["text"].split("\n")):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = code_line
                run.font.size = Pt(14)
                run.font.name = FONT_MONO
                run.font.color.rgb = COLOR_CODE_FG
            current_top += code_height + Inches(0.2)

        elif block["type"] == "quote":
            txBox = slide.shapes.add_textbox(
                Inches(1.4), current_top, Inches(10.5), Inches(0.8)
            )
            # Add left accent bar
            bar = slide.shapes.add_shape(
                1, MARGIN_LEFT, current_top, Pt(4), Inches(0.7)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_ACCENT
            bar.line.fill.background()

            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _add_formatted_runs(p, block["text"], base_size=18,
                                base_color=COLOR_QUOTE, base_italic=True)
            current_top += Inches(0.9)

        elif block["type"] == "table":
            _render_table(slide, block["lines"], current_top)
            current_top += Inches(len(block["lines"]) * 0.4 + 0.2)

        elif block["type"] == "paragraph":
            txBox = slide.shapes.add_textbox(
                MARGIN_LEFT, current_top, CONTENT_WIDTH, Inches(0.8)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.space_after = Pt(8)
            _add_formatted_runs(p, block["text"], base_size=18)
            current_top += Inches(0.5)

        elif block["type"] in ("h2", "h3"):
            txBox = slide.shapes.add_textbox(
                MARGIN_LEFT, current_top, CONTENT_WIDTH, Inches(0.6)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            size = 24 if block["type"] == "h2" else 20
            _add_formatted_runs(p, block["text"], base_size=size,
                                base_color=COLOR_HEADING, base_bold=True)
            current_top += Inches(0.6)


def _render_table(slide, table_lines: list[str], top: float):
    """Render a markdown table as a python-pptx table."""
    # Parse table
    rows_raw: list[list[str]] = []
    for line in table_lines:
        cells = [c.strip() for c in line.strip("|").split("|")]
        # Skip separator rows (---|----|---)
        if all(re.match(r'^[-:]+$', c) for c in cells):
            continue
        rows_raw.append(cells)

    if not rows_raw:
        return

    n_rows = len(rows_raw)
    n_cols = max(len(r) for r in rows_raw)

    # Ensure all rows have same number of columns
    for r in rows_raw:
        while len(r) < n_cols:
            r.append("")

    col_width = int(CONTENT_WIDTH / n_cols)
    row_height = Inches(0.4)

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        MARGIN_LEFT, top,
        CONTENT_WIDTH, Inches(n_rows * 0.4)
    ).table

    for i, row in enumerate(rows_raw):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            if i == 0:
                # Header row styling
                _add_formatted_runs(p, cell_text, base_size=13,
                                    base_color=COLOR_WHITE, base_bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_HEADER_BG
            else:
                _add_formatted_runs(p, cell_text, base_size=13, base_color=COLOR_BODY)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_TABLE_ALT_BG

            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)


# ── Main conversion ───────────────────────────────────────────────────────
def convert(input_path: str, output_path: str, title: str | None = None):
    """Convert a markdown slide file to .pptx."""
    md_content = Path(input_path).read_text(encoding="utf-8")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Split into slides
    raw_slides = md_content.split("\n---\n")
    slides = [s.strip() for s in raw_slides if s.strip()]

    if not slides:
        print("No slides found in input file.", file=sys.stderr)
        sys.exit(1)

    for idx, slide_md in enumerate(slides):
        blocks = _parse_slide_blocks(slide_md)
        if not blocks:
            continue

        # First slide with only h1 (+ optional subtitle) → title slide
        if idx == 0 and _is_title_slide(blocks):
            _render_title_slide(prs, blocks)
        else:
            _render_content_slide(prs, blocks)

    prs.save(output_path)
    print(f"Saved {len(slides)} slides to {output_path}")


def main():
    parser = argparse.ArgumentParser(description="Convert markdown slides to PPTX")
    parser.add_argument("input", help="Input markdown file")
    parser.add_argument("-o", "--output", default=None, help="Output .pptx file")
    parser.add_argument("--title", default=None, help="Deck title (optional)")
    args = parser.parse_args()

    output = args.output
    if output is None:
        output = Path(args.input).with_suffix(".pptx").name

    convert(args.input, output, args.title)


if __name__ == "__main__":
    main()
